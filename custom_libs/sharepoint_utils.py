import io
import os
from datetime import datetime, timedelta, timezone
from typing import Any, Dict, List, Optional, Union
from pptx import Presentation
import PyPDF2
import msal
import requests
from docx import Document as DocxDocument
from dotenv import load_dotenv
from custom_libs.custom_logging import get_logger
import random

logger = get_logger()

class SharePointUtils:
    """This class facilitates the extraction of data from SharePoint using Microsoft Graph API.
    It supports authentication and data retrieval from SharePoint sites, lists, and libraries.
    """

    def __init__(
        self,
        tenant_id: Optional[str] = None,
        client_id: Optional[str] = None,
        client_secret: Optional[str] = None,
        graph_uri: str = "https://graph.microsoft.com",
        authority_template: str = "https://login.microsoftonline.com/{tenant_id}",
    ):
        """
        Initialize the SharePointDataExtractor class with optional environment variables.

        :param tenant_id: Tenant ID for Microsoft 365.
        :param client_id: Client ID for the application registered in Azure AD.
        :param client_secret: Client secret for the application registered in Azure AD.
        :param graph_uri: URI for Microsoft Graph API.
        :param authority_template: Template for authority URL used in authentication.
        """
        self.tenant_id = tenant_id
        self.client_id = client_id
        self.client_secret = client_secret
        self.graph_uri = graph_uri
        self.authority = (
            authority_template.format(tenant_id=tenant_id) if tenant_id else None
        )
        self.scope = ["https://graph.microsoft.com/.default"]
        self.access_token = None

    def loadEnvFile(self):
        """
        Loads required environment variables for the application from a .env file.

        This method should be called explicitly if environment variables are to be loaded from a .env file.
        """
        load_dotenv()

        self.tenant_id = os.getenv("AZURE_TENANT_ID")
        self.client_id = os.getenv("AZURE_CLIENT_ID")
        self.client_secret = os.getenv("AZURE_CLIENT_SECRET")
        self.authority = f"https://login.microsoftonline.com/{self.tenant_id}"
        self.app_name = os.getenv('APP_NAME')
        self.purview_acct_name = os.getenv('PURVIEW_ACCOUNT_NAME')
        self.purview_endpoint_url = os.getenv('PURVIEW_ENDPOINT_URL')
        self.purview_token_url = os.getenv('PURVIEW_TOKEN_URL')
        self.sharepoint_site_domain = os.getenv('SITE_DOMAIN')
        self.sharepoint_site_name = os.getenv('SITE_NAME')
        self.azure_openai_api_key = os.getenv('AZURE_OPENAI_API_KEY')
        self.azure_openai_endpoint = os.getenv('AZURE_OPENAI_ENDPOINT')
        self.azure_openai_deployment = os.getenv('AZURE_OPENAI_DEPLOYMENT_NAME')
        self.azure_openai_llm_model = os.getenv('AZURE_OPENAI_LLM_MODEL')
        self.azure_openai_api_version = os.getenv("AZURE_OPENAI_API_VERSION")

        # Check for any missing required environment variables
        missing_vars = [
            var_name
            for var_name, var in [
                ("AZURE_TENANT_ID", self.tenant_id),
                ("AZURE_CLIENT_ID", self.client_id),
                ("AZURE_CLIENT_SECRET", self.client_secret),
                ('APP_NAME', self.app_name),
                ('PURVIEW_ACCOUNT_NAME', self.purview_acct_name),
                ('PURVIEW_ENDPOINT_URL', self.purview_endpoint_url),
                ('PURVIEW_TOKEN_URL', self.purview_token_url),
                ('SITE_DOMAIN', self.sharepoint_site_domain),
                ('SITE_NAME', self.sharepoint_site_name),
                ('AZURE_OPENAI_API_KEY', self.azure_openai_api_key),
                ('AZURE_OPENAI_ENDPOINT', self.azure_openai_endpoint),
                ('AZURE_OPENAI_DEPLOYMENT_NAME', self.azure_openai_deployment),
                ('AZURE_OPENAI_LLM_MODEL', self.azure_openai_llm_model),
                ('AZURE_OPENAI_API_VERSION', self.azure_openai_api_version),
            ]
            if not var
        ]

        if missing_vars:
            raise EnvironmentError(
                f"Missing required environment variables: {', '.join(missing_vars)}"
            )

        environmentVariables = [
            "AZURE_TENANT_ID", 
            "AZURE_CLIENT_ID", 
            "AZURE_CLIENT_SECRET",
            "APP_NAME",
            "PURVIEW_ACCOUNT_NAME",
            "PURVIEW_ENDPOINT_URL",
            "PURVIEW_TOKEN_URL"
            "SITE_DOMAIN",
            "SITE_NAME",
            "AZURE_OPENAI_API_KEY",
            "AZURE_OPENAI_ENDPOINT",
            "AZURE_OPENAI_DEPLOYMENT_NAME",
            "AZURE_OPENAI_LLM_MODEL",
            "AZURE_OPENAI_API_VERSION"
        ]
        # Log the success of loading each environment variable
        loaded_vars = [
            var_name
            for var_name in environmentVariables
            if var_name not in missing_vars
        ]
        if loaded_vars:
            logger.info(
                f"Successfully loaded environment variables: {', '.join(loaded_vars)}"
            )

    def msgraph_auth(
        self,
        client_id: Optional[str] = None,
        client_secret: Optional[str] = None,
        authority: Optional[str] = None,
    ):
        """
        Authenticate with Microsoft Graph using MSAL for Python.
        """
        # Use provided parameters or fall back to instance attributes
        client_id = client_id or self.client_id
        client_secret = client_secret or self.client_secret
        authority = authority or self.authority

        # Check if all necessary credentials are provided
        if not all([client_id, client_secret, authority]):
            raise ValueError("Missing required authentication credentials.")

        app = msal.ConfidentialClientApplication(
            client_id=client_id, authority=authority, client_credential=client_secret
        )

        try:
            # Attempt to acquire token
            access_token = app.acquire_token_silent(self.scope, account=None)
            if not access_token:
                access_token = app.acquire_token_for_client(scopes=self.scope)
                if "access_token" in access_token:
                    logger.info("New access token retrieved.")
                else:
                    logger.error("Error acquiring authorization token.")
                    return None
            else:
                logger.info("Token retrieved from MSAL Cache.")

            # Store the access token in the instance
            self.access_token = access_token["access_token"]
            return self.access_token

        except Exception as err:
            logger.error(f"Error in msgraph_auth: {err}")
            raise

    @staticmethod
    def _format_url(site_id: str, drive_id: str, folder_path: str = None) -> str:
        """
        Formats the URL for accessing a nested site drive in Microsoft Graph.

        :param site_id: The site ID in Microsoft Graph.
        :param drive_id: The drive ID in Microsoft Graph.
        :param folder_path: path to the folder within the drive, can include subfolders.
            The format should follow '/folder/subfolder1/subfolder2/'. For example,
            '/test/test1/test2/' to access nested folders.
        :return: The formatted URL.
        """
        folder_path_formatted = folder_path.rstrip("/")
        return f"https://graph.microsoft.com/v1.0/sites/{site_id}/drives/{drive_id}/root:{folder_path_formatted}:/"

    def _make_ms_graph_request(
        self, url: str, access_token: Optional[str] = None
    ) -> Dict:
        """
        Make a request to the Microsoft Graph API.

        :param url: The URL for the Microsoft Graph API endpoint.
        :param access_token: Optional; The access token for Microsoft Graph API authentication. If not provided, uses the instance's stored token.
        :return: The JSON response from the Microsoft Graph API.
        :raises Exception: If there's an HTTP error or other issues in making the request.
        """
        access_token = access_token or self.access_token
        if not access_token:
            raise ValueError("Access token is required for making API requests.")

        headers = {"Authorization": f"Bearer {access_token}"}
        try:
            response = requests.get(url, headers=headers)
            response.raise_for_status()
            return response.json()
        except requests.exceptions.HTTPError as err:
            logger.error(f"HTTP Error: {err}")
            raise
        except Exception as err:
            logger.error(f"Error in _make_ms_graph_request: {err}")
            raise

    def get_site_id(
        self, site_domain: str, site_name: str, access_token: Optional[str] = None
    ) -> Optional[str]:
        """
        Get the Site ID from Microsoft Graph API.
        """
        endpoint = (
            f"https://graph.microsoft.com/v1.0/sites/{site_domain}:/sites/{site_name}:/"
        )
        access_token = access_token or self.access_token

        try:
            logger.info("Getting the Site ID...")
            result = self._make_ms_graph_request(endpoint, access_token)
            site_id = result.get("id")
            if site_id:
                logger.info(f"Site ID retrieved: {site_id}")
                return site_id
        except Exception as err:
            logger.error(f"Error retrieving Site ID: {err}")
            return None

    def get_drive_id(self, site_id: str, access_token: Optional[str] = None) -> str:
        """
        Get the drive ID from a Microsoft Graph site.
        """
        url = f"https://graph.microsoft.com/v1.0/sites/{site_id}/drive"

        access_token = access_token or self.access_token

        try:
            json_response = self._make_ms_graph_request(url, access_token)
            drive_id = json_response.get("id")
            logger.info(f"Successfully retrieved drive ID: {drive_id}")
            return drive_id
        except Exception as err:
            logger.error(f"Error in get_drive_id: {err}")
            raise

    def get_files_in_site(
        self,
        site_id: str,
        drive_id: str,
        folder_path: Optional[str] = None,
        access_token: Optional[str] = None,
        minutes_ago: Optional[int] = None,
        file_formats: Optional[List[str]] = None,
    ) -> List[Dict]:
        """
        Get a list of files in a site's drive, optionally filtered by creation or last modification time and file formats.

        :param site_id: The site ID in Microsoft Graph.
        :param drive_id: The drive ID in Microsoft Graph.
        :param folder_path: Path to the folder within the drive, can include subfolders.
                The format should follow '/folder/subfolder1/subfolder2/'.For example,
                '/test/test1/test2/' to access nested folders.
        :param access_token: The access token for Microsoft Graph API authentication. If not provided, it will be fetched from self.
        :param minutes_ago: Optional integer to filter files created or updated within the specified number of minutes from now.
        :param file_formats: List of desired file formats.
        :return: A list of file details.
        :raises Exception: If there's an error in fetching file details.
        """
        if access_token is None:
            access_token = self.access_token

        # Construct the URL based on whether a folder path is provided
        if folder_path:
            url = self._format_url(site_id, drive_id, folder_path) + "children"
        else:
            url = f"https://graph.microsoft.com/v1.0/sites/{site_id}/drives/{drive_id}/root/children"

        try:
            logger.info("Making request to Microsoft Graph API")
            json_response = self._make_ms_graph_request(url, access_token)
            files = json_response["value"]
            logger.info("Received response from Microsoft Graph API")

            time_limit = (
                datetime.now(timezone.utc) - timedelta(minutes=minutes_ago)
                if minutes_ago is not None
                else None
            )

            filtered_files = [
                file
                for file in files
                    if (
                        (
                            time_limit is None
                            or datetime.fromisoformat(
                                file["fileSystemInfo"]["createdDateTime"].rstrip("Z")
                            ).replace(tzinfo=timezone.utc)
                            >= time_limit
                            or datetime.fromisoformat(
                                file["fileSystemInfo"]["lastModifiedDateTime"].rstrip("Z")
                            ).replace(tzinfo=timezone.utc)
                            >= time_limit
                        )
                        and (
                            not file_formats
                            or any(file["name"].endswith(f".{fmt}") for fmt in file_formats)
                        )
                    )
            ]

            return filtered_files
        except Exception as err:
            logger.error(f"Error in get_files_in_site: {err}")
            raise

    def get_file_permissions(
        self, site_id: str, item_id: str, access_token: Optional[str] = None
    ) -> List[Dict]:
        """
        Get the permissions of a file in a site.

        :param site_id: The site ID in Microsoft Graph.
        :param item_id: The item ID of the file in Microsoft Graph.
        :param access_token: The access token for Microsoft Graph API authentication. If not provided, it will be fetched from self.
        :return: A list of permission details.
        :raises Exception: If there's an error in fetching permission details.
        """
        if access_token is None:
            access_token = self.access_token

        url = f"https://graph.microsoft.com/v1.0/sites/{site_id}/drive/items/{item_id}/permissions"

        try:
            json_response = self._make_ms_graph_request(url, access_token)
            return json_response["value"]
        except Exception as err:
            logger.error(f"Error in get_file_permissions: {err}")
            raise

    @staticmethod
    def get_read_access_entities(permissions):
        """
        Extracts user IDs and group names of entities with read access from the given permissions data.

        :param permissions: List of permission dictionaries.
        :return: List of entities (user IDs and group names/IDs) with read access.
        """
        read_access_entities = []

        for permission in permissions:
            if not isinstance(permission, dict) or "roles" not in permission:
                continue

            if "read" in permission.get("roles", []):
                # Process grantedToIdentitiesV2 for individual users
                identities_v2 = permission.get("grantedToIdentitiesV2", [])
                for identity in identities_v2:
                    user = identity.get("user", {})
                    user_id = user.get("id")
                    if user_id and user_id not in read_access_entities:
                        read_access_entities.append(user_id)

                # Process grantedToIdentities for individual users
                identities = permission.get("grantedToIdentities", [])
                for identity in identities:
                    user = identity.get("user", {})
                    user_id = user.get("id")
                    if user_id and user_id not in read_access_entities:
                        read_access_entities.append(user_id)

                # Process grantedToV2 for groups
                groups = permission.get("grantedToV2", {}).get("siteGroup", {})
                group_name = groups.get(
                    "displayName"
                )  # or groups.get('id') for group ID
                if group_name and group_name not in read_access_entities:
                    read_access_entities.append(group_name)

        return read_access_entities

    def get_file_content_bytes(
        self,
        site_id: str,
        drive_id: str,
        folder_path: Optional[str],
        file_name: str,
        access_token: Optional[str] = None,
    ) -> Optional[bytes]:
        """
        Retrieve the content of a file as bytes from a specific site drive.

        :param site_id: The site ID in Microsoft Graph.
        :param drive_id: The drive ID in Microsoft Graph.
        :param folder_path: Path to the folder within the drive, can include subfolders.
        :param file_name: The name of the file.
        :param specific_file: Specific file name, if different from file_name.
        :param access_token: The access token for Microsoft Graph API authentication.
        :return: Bytes content of the file or None if there's an error.
        """
        if access_token is None:
            access_token = self.access_token

        folder_path_formatted = folder_path.rstrip("/") if folder_path else ""
        endpoint = f"https://graph.microsoft.com/v1.0/sites/{site_id}/drives/{drive_id}/root:{folder_path_formatted}/{file_name}:/content"

        try:
            response = requests.get(
                endpoint, headers={"Authorization": "Bearer " + access_token}
            )
            if response.status_code != 200:
                logger.error(
                    f"Failed to retrieve file content. Status code: {response.status_code}, Response: {response.text}"
                )
                return None
            return response.content
        except requests.exceptions.RequestException as req_err:
            logger.error(f"Request error: {req_err}")
            return None

    def process_and_retrieve_docx_content(
        self,
        site_id: str,
        drive_id: str,
        folder_path: Optional[str],
        file_name: str,
        access_token: Optional[str] = None,
    ) -> Optional[str]:
        """
        Process the content of a .docx file and extract its text.

        :param site_id: The site ID in Microsoft Graph.
        :param drive_id: The drive ID in Microsoft Graph.
        :param folder_path: Path to the folder within the drive, can include subfolders.
        :param file_name: The name of the .docx file.
        :param access_token: The access token for Microsoft Graph API authentication.
        :return: Text content of the .docx file or None if there's an error.
        """
        file_bytes = self.get_file_content_bytes(
            site_id, drive_id, folder_path, file_name, access_token
        )
        if file_bytes is None:
            return None

        if not file_name.endswith(".docx"):
            logger.error(f"File {file_name} is not a .docx file.")
            return None

        try:
            document = DocxDocument(io.BytesIO(file_bytes))
            return "\n".join([paragraph.text for paragraph in document.paragraphs])
        except Exception as err:
            logger.error(f"Error processing document: {err}")
            return None

    def process_and_retrieve_pptx_content(
        self,
        site_id: str,
        drive_id: str,
        folder_path: Optional[str],
        file_name: str,
        access_token: Optional[str] = None,
    ) -> Optional[str]:
        """
        Process the content of a .docx file and extract its text.

        :param site_id: The site ID in Microsoft Graph.
        :param drive_id: The drive ID in Microsoft Graph.
        :param folder_path: Path to the folder within the drive, can include subfolders.
        :param file_name: The name of the .docx file.
        :param access_token: The access token for Microsoft Graph API authentication.
        :return: Text content of the .docx file or None if there's an error.
        """
        file_bytes = self.get_file_content_bytes(
            site_id, drive_id, folder_path, file_name, access_token
        )
        if file_bytes is None:
            return None
        
        if not file_name.endswith(".pptx"):
            logger.error(f"File {file_name} is not a .pptx file.")
            return None
        try:
            document = ""
            prs = Presentation(io.BytesIO(file_bytes))
            # prs = Presentation(file_bytes)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        document += shape.text
            return document
        except Exception as err:
            logger.error(f"Error processing document: {err}")
            return None

    def extract_text_from_pdf_bytes(
            self,
            pdf_bytes: bytes
    ) -> Optional[str]:
        """
        Extracts text from a PDF file provided as a bytes object.

        :param pdf_bytes: Bytes object containing the PDF file data.
        :return: Extracted text from the PDF as a string, or None if extraction fails.
        """
        try:
            with io.BytesIO(pdf_bytes) as pdf_stream:
                pdf_reader = PyPDF2.PdfReader(pdf_stream)
                text = []
                for page_num in range(len(pdf_reader.pages)):
                    page = pdf_reader.pages[page_num]
                    text.append(page.extract_text())

                extracted_text = "\n".join(text)
                logger.info("Text extraction from PDF bytes was successful.")
                return extracted_text
        except Exception as e:
            logger.error(f"An unexpected error occurred during PDF text extraction: {e}")

        return None

    def process_and_retrieve_pdf_content(
        self,
        site_id: str,
        drive_id: str,
        folder_path: Optional[str],
        file_name: str,
        access_token: Optional[str] = None,
    ) -> Optional[str]:
        """
        Process the content of a .docx file and extract its text.

        :param site_id: The site ID in Microsoft Graph.
        :param drive_id: The drive ID in Microsoft Graph.
        :param folder_path: Path to the folder within the drive, can include subfolders.
        :param file_name: The name of the .docx file.
        :param specific_file: Specific .docx file name, if different from file_name.
        :param access_token: The access token for Microsoft Graph API authentication.
        :return: Text content of the .docx file or None if there's an error.
        """
        file_bytes = self.get_file_content_bytes(
            site_id, drive_id, folder_path, file_name, access_token
        )
        if file_bytes is None:
            return None

        if not file_name.endswith(".pdf"):
            logger.error(f"File {file_name} is not a .pdf file.")
            return None

        try:
            document = self.extract_text_from_pdf_bytes(file_bytes)
            return document
        except Exception as err:
            logger.error(f"Error processing document: {err}")
            return None

    @staticmethod
    def _extract_file_metadata(
        file_data: Dict[str, Any]
    ) -> Dict[str, Optional[Union[str, datetime]]]:
        """
        Extracts specific information from the file data.

        This function takes a dictionary containing file data and returns a new dictionary
        with specific fields: 'webUrl', 'size', 'createdBy', 'createdDateTime',
        'lastModifiedDateTime', and 'lastModifiedBy'.

        Args:
            file_data (Dict[str, Any]): The original file data.

        Returns:
            Dict[str, Optional[Union[str, datetime]]]: A dictionary with the extracted file information.
            If a field is not present in the file data, the function will return None for that field.
        """

        def format_date(date_str):
            # Append 'Z' if it's missing to indicate UTC timezone
            return date_str if date_str.endswith("Z") else f"{date_str}Z"

        return {
            "id": file_data.get("id"),
            "webUrl": file_data.get("webUrl"),
            "size": file_data.get("size"),
            "createdBy": file_data.get("createdBy", {})
            .get("user", {})
            .get("displayName"),
            "createdDateTime": format_date(
                file_data.get("fileSystemInfo", {}).get("createdDateTime", "")
            )
            if file_data.get("fileSystemInfo", {}).get("createdDateTime")
            else None,
            "lastModifiedDateTime": format_date(
                file_data.get("fileSystemInfo", {}).get("lastModifiedDateTime", "")
            )
            if file_data.get("fileSystemInfo", {}).get("lastModifiedDateTime")
            else None,
            "lastModifiedBy": file_data.get("lastModifiedBy", {})
            .get("user", {})
            .get("displayName"),
        }

    def getSharepointFileContent(
        self,
        site_domain: str,
        site_name: str,
        folder_path: Optional[str] = None,
        file_names: Optional[Union[str, List[str]]] = None,
        minutes_ago: Optional[int] = None,
        file_formats: Optional[List[str]] = None,
    ) -> Dict[str, Dict[str, Optional[str]]]:
        """
        Retrieve contents of files from a specified SharePoint location, optionally filtering by last modification time and file formats.

        :param site_domain: The domain of the site in Microsoft Graph.
        :param site_name: The name of the site in Microsoft Graph.
        :param folder_path: Path to the folder within the drive, can include subfolders like 'test1/test2'.
        :param file_names: Optional; the name or names of specific files to retrieve. If provided, only these files' content will be fetched.
        :param minutes_ago: Optional; filter for files modified within the specified number of minutes.
        :param file_formats: Optional; list of desired file formats to include.
        :return: Dictionary with file names as keys and a dictionary containing their content, location, and users_by_role as values.
        """
        if self._are_required_variables_missing():
            return None

        site_id, drive_id = self._get_site_and_drive_ids(site_domain, site_name)

        if not site_id or not drive_id:
            return None

        files = self._get_files(
            site_id, drive_id, folder_path, minutes_ago, file_formats
        )
        if not files:
            logger.error("No files found in the site's drive")
            return None

        return self._process_files(
            # Israel
            # site_id, drive_id, folder_path, file_names, files, file_formats
            site_domain, site_name, site_id, drive_id, folder_path, file_names, files, file_formats
        )

    def listSharepointFiles(
        self,
        site_domain: str,
        site_name: str,
        folder_path: Optional[str] = None,
        file_names: Optional[Union[str, List[str]]] = None,
        minutes_ago: Optional[int] = None,
        file_formats: Optional[List[str]] = None,
    ) -> Dict[str, Dict[str, Optional[str]]]:
        """
        Retrieve contents of files from a specified SharePoint location, optionally filtering by last modification time and file formats.

        :param site_domain: The domain of the site in Microsoft Graph.
        :param site_name: The name of the site in Microsoft Graph.
        :param folder_path: Path to the folder within the drive, can include subfolders like 'test1/test2'.
        :param file_names: Optional; the name or names of specific files to retrieve. If provided, only these files' content will be fetched.
        :param minutes_ago: Optional; filter for files modified within the specified number of minutes.
        :param file_formats: Optional; list of desired file formats to include.
        :return: Dictionary with file names as keys and a dictionary containing their content, location, and users_by_role as values.

        Added by Israel Oros
        """
        if self._are_required_variables_missing():
            return None

        site_id, drive_id = self._get_site_and_drive_ids(site_domain, site_name)

        if not site_id or not drive_id:
            return None

        files = self._get_files(
            site_id, drive_id, folder_path, minutes_ago, file_formats
        )
        if not files:
            logger.error("No files found in the site's drive")
            return None

        return files

    def sharepointFileSampleList(
            self,
            input_file_list: List[Dict[str, str]], 
            file_sample_size: int
    ) -> List[str]:
        """
        Generates a list of files to process from an input list of files.

        Args:
            input_file_list (list): A list of dictionaries containing file information.
                Each dictionary should have a 'name' key representing the file name.
            file_sample_size (int): The desired sample size of files to process.

        Returns:
            list: A list of file names to process. If the number of files found is greater
            than the sample size, a random sample of files is returned. Otherwise, all files
            in the input list are included.

        Example:
            input_files = [{'name': 'file1.txt'}, {'name': 'file2.txt'}, {'name': 'file3.txt'}]
            sample_size = 2
            result = sharepointFileSampleList(input_files, sample_size)
            # Example output: ['file1.txt', 'file3.txt']
        """
        num_files_found = len(input_file_list)
        files_to_process = []

        if num_files_found > file_sample_size:
            # Randomly sample files if more files are found than the desired sample size
            # for ind in random.sample(range(num_files_found), file_sample_size):
            #     files_to_process.append(input_file_list[ind]['name'])
            files_to_process = [input_file_list[ind]['name'] for ind in random.sample(range(num_files_found), file_sample_size)]
        else:
            # Include all files if the number of files found is less than or equal to the sample size
            # for ind in input_file_list:
            #     files_to_process.append(ind['name'])
            files_to_process = [file['name'] for file in input_file_list]

        return files_to_process

    def _are_required_variables_missing(self) -> bool:
        """
        Checks if any of the required instance variables for SharePointDataExtractor are missing.

        This function checks the following instance variables: 'tenant_id', 'client_id',
        'client_secret', 'graph_uri', and 'authority'. If any of these variables are not set,
        the function logs an error message and returns True.

        :return: True if any of the required instance variables are missing, False otherwise.
        """
        required_vars = {
            "tenant_id": self.tenant_id,
            "client_id": self.client_id,
            "client_secret": self.client_secret,
            "graph_uri": self.graph_uri,
            "authority": self.authority,
        }
        missing_vars = [var_name for var_name, var in required_vars.items() if not var]
        if missing_vars:
            logger.error(
                f"Required instance variables for SharePointDataExtractor are not set: {', '.join(missing_vars)}. Please load loadEnvFile or set them manually."
            )
            return True
        return False

    def _get_site_and_drive_ids(
        self, site_domain: str, site_name: str
    ) -> (Optional[str], Optional[str]):
        """
        Retrieves the site ID and drive ID for a given site domain and site name.

        :param site_domain: The domain of the site.
        :param site_name: The name of the site.
        :return: A tuple containing the site ID and drive ID, or (None, None) if either ID could not be retrieved.
        """
        site_id = self.get_site_id(site_domain, site_name)
        if not site_id:
            logger.error("Failed to retrieve site_id")
            return None, None

        drive_id = self.get_drive_id(site_id)
        if not drive_id:
            logger.error("Failed to retrieve drive ID")
            return None, None

        return site_id, drive_id

    def _get_files(
        self,
        site_id: str,
        drive_id: str,
        folder_path: Optional[str],
        minutes_ago: Optional[int],
        file_formats: Optional[List[str]],
    ) -> List[Dict]:
        """
        Retrieves the files in a site drive.

        :param site_id: The site ID in Microsoft Graph.
        :param drive_id: The drive ID in Microsoft Graph.
        :param folder_path: Optional path to the folder within the drive, can include subfolders.
        :param minutes_ago: Optional integer to filter files created or updated within the specified number of minutes from now.
        :param file_formats: List of desired file formats.
        :param specific_file: Optional specific file to retrieve.
        :return: A list of file details.
        """
        files = self.get_files_in_site(
            site_id=site_id,
            drive_id=drive_id,
            folder_path=folder_path,
            minutes_ago=minutes_ago,
            file_formats=file_formats,
        )
        return files

    def _process_files(
        self,
        # Israel#######
        site_domain: str, 
        site_name: str,
        ##########
        site_id: str,
        drive_id: str,
        folder_path: Optional[str],
        file_names: Optional[Union[str, List[str]]],
        files: List[Dict],
        file_formats: Optional[List[str]],
    ) -> List[Dict[str, Optional[str]]]:
        """Processes the files in a site drive.
        # Israel
        :param site_domain: The domain for the SPO site.
        :param site_name: The name of the SPO site.
        ########
        :param site_id: The site ID in Microsoft Graph.
        :param drive_id: The drive ID in Microsoft Graph.
        :param folder_path: Optional path to the folder within the drive, can include subfolders.
        :param file_names: The name(s) of specific files to filter. Can be a string or a list of strings.
        :param files: List of files to process.
        :param file_formats: List of desired file formats.
        :return: A list of dictionaries, each mapping file names to their content and metadata.
        """
        file_contents = []

        # Handle both string and list for file_names
        if isinstance(file_names, str):
            file_names = [file_names]

        # Filter files based on the given file_names
        if file_names:
            files = [file for file in files if file.get("name") in file_names]
            if len(files) == 0:
                logger.error("No matching files found")
                return []

        for file in files:
            print(f"Processing File {file}")
            file_name = file.get("name")
            if file_name and self._is_file_format_valid(file_name, file_formats):
                metadata = self._extract_file_metadata(file)
                content = self._retrieve_file_content(
                    site_id, drive_id, folder_path, file_name
                )
                users_by_role = self.get_read_access_entities(
                    self.get_file_permissions(site_id, file["id"])
                )
                file_content = {
                    "content": content,
                    # Israel
                    # **self._format_metadata(metadata, file_name, users_by_role),
                    **self._format_metadata(metadata, file_name, users_by_role, site_domain, site_name, folder_path),
                }
                file_contents.append(file_content)

        return file_contents

    def _is_file_format_valid(
        self, file_name: str, file_formats: Optional[List[str]]
    ) -> bool:
        """
        Checks if the format of a file is valid.

        :param file_name: The name of the file.
        :param file_formats: List of desired file formats.
        :return: True if the file format is valid, False otherwise.
        """
        return "." in file_name and (
            not file_formats
            or any(file_name.endswith(f".{fmt}") for fmt in file_formats)
        )

    def _retrieve_file_content(
        self, site_id: str, drive_id: str, folder_path: Optional[str], file_name: str
    ) -> Optional[str]:
        """
        Retrieve the content of a specific file from SharePoint.

        :param site_id: SharePoint site ID.
        :param drive_id: SharePoint drive ID.
        :param folder_path: Path to the folder containing the file.
        :param file_name: Name of the file to retrieve.
        :return: Content of the file as a string, or None if retrieval fails.
        """
        if file_name.endswith(".docx"):
            return self.process_and_retrieve_docx_content(
                site_id, drive_id, folder_path, file_name
            )
        elif file_name.endswith(".pdf"):
            return self.process_and_retrieve_pdf_content(
                site_id, drive_id, folder_path, file_name
            )
        elif file_name.endswith(".pptx"):
            return self.process_and_retrieve_pptx_content(
                site_id, drive_id, folder_path, file_name
            )         
        # Add other file type processing as needed
        return None

    def _retrieve_fs_file_content(
        self, folder_path: str, file_name: str
    ) -> Optional[str]:
        """
        Retrieve the content of a specific file from SharePoint.

        :param site_id: SharePoint site ID.
        :param drive_id: SharePoint drive ID.
        :param folder_path: Path to the folder containing the file.
        :param file_name: Name of the file to retrieve.
        :return: Content of the file as a string, or None if retrieval fails.
        """
        if file_name.endswith(".docx"):
            return self.process_and_retrieve_docx_content(
                site_id, drive_id, folder_path, file_name
            )
        elif file_name.endswith(".pdf"):
            return self.process_and_retrieve_pdf_content(
                site_id, drive_id, folder_path, file_name
            )
        elif file_name.endswith(".pptx"):
            return self.process_and_retrieve_pptx_content(
                site_id, drive_id, folder_path, file_name
            )         

        # Add other file type processing as needed
        return None
    
    def _format_metadata(
        self,
        metadata: Dict,
        file_name: str,
        users_by_role: Dict,
        # Israel
        site_domain: str,
        site_name: str,
        folder_path: str,
        #######
    ) -> Dict:
        """
        Format and return file metadata.

        :param metadata: Dictionary of file metadata.
        :param file_name: Name of the file.
        :param users_by_role: Dictionary of users grouped by their role.
        :return: Formatted metadata as a dictionary.
        """
        formatted_metadata = {
            "id": metadata["id"],
            "source": metadata["webUrl"],
            "name": file_name,
            "size": metadata["size"],
            "created_by": metadata["createdBy"],
            "created_datetime": metadata["createdDateTime"],
            "last_modified_datetime": metadata["lastModifiedDateTime"],
            "last_modified_by": metadata["lastModifiedBy"],
            "read_access_entity": users_by_role,
            # Israel
            "parentObject": f"{site_domain}/{site_name}{folder_path}",
            #######
            "typedef": "SharePoint"
        }
        return formatted_metadata