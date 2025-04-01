import requests
import json
import re
import tiktoken
from azure.ai.inference.models import SystemMessage, UserMessage
from pyapacheatlas.core import AtlasClassification, AtlasEntity
from datetime import datetime
import os
import random
from pathlib import Path
import uuid
import PyPDF2
from pptx import Presentation
from docx import Document as DocxDocument
from custom_libs.custom_logging import get_logger

logger = get_logger()

def filesystemFileSampleList(input_file_list, file_sample_size,file_system_path,):
    """
    From a list of files identified, generates a subset of files based on the given sample
    size, extracts the contents from each file, and returns a list.

    Args:
        input_file_list (list): A list of file paths (as strings).
        file_sample_size (int): The desired sample size of files to process.
        file_system_path (str): The file system path to be used as a parent object.

    Returns:
        list: A list of dictionaries containing file information. Each dictionary includes:
            - 'id': A unique identifier (UUID) for the file.
            - 'name': The name of the file.
            - 'created_datetime': The creation timestamp of the file.
            - 'created_by': The user ID of the file creator.
            - 'size': The file size in bytes.
            - 'last_modified_datetime': The last modification timestamp of the file.
            - 'last_modified_by': The user ID of the last modifier.
            - 'source': The original file path.

    Example:
        input_files = ['/path/to/file1.txt', '/path/to/file2.txt', '/path/to/file3.txt']
        sample_size = 2
        result = filesystemFileSampleList(input_files, sample_size)
        # Example output (for illustration purposes):
        # [{'id': '12345', 'name': 'file1.txt', 'created_datetime': ..., 'created_by': 'user1', ...},
        #  {'id': '67890', 'name': 'file3.txt', 'created_datetime': ..., 'created_by': 'user2', ...}]
    """
    num_files_found = len(input_file_list)
    files_to_process = []
    content = ''

    if num_files_found > file_sample_size:
        # Randomly sample files if more files are found than the desired sample size
        for item in random.sample(range(num_files_found), file_sample_size):
            fileFullPath = Path(input_file_list[item])
            fileName = fileFullPath.name
            print(f"Processing {fileName}")
            if fileName.endswith(".docx"):
                try:
                    document = DocxDocument(fileFullPath)
                    content = "\n".join([paragraph.text for paragraph in document.paragraphs])
                except Exception as err:
                    logger.error(f"Error processing document: {err}")
            elif fileName.endswith(".pdf"):
                try:
                    content = extractPDFContent(fileFullPath)
                except Exception as err:
                    logger.error(f"Error processing document: {err}")
            elif fileName.endswith(".pptx"):
                try:
                    content = extractPPTXContent(fileFullPath)
                except Exception as err:
                    logger.error(f"Error processing document: {err}")
            # Create a dictionary with file information
            file_info = {
                'id': str(uuid.uuid4()),
                'name': fileName,
                'created_datetime': datetime.fromtimestamp(fileFullPath.stat().st_ctime),
                'created_by': str(fileFullPath.stat().st_uid),
                'size': fileFullPath.stat().st_size,
                'last_modified_datetime': datetime.fromtimestamp(fileFullPath.stat().st_mtime),
                'last_modified_by': str(fileFullPath.stat().st_uid),
                'source': input_file_list[item],
                # Israel
                "parentObject": f"{file_system_path}",
                ########
                'typedef': 'FileSystem',
                'content': content,
            }
            files_to_process.append(file_info)
    else:
        # Include all files if the number of files found is less than or equal to the sample size
        for item in input_file_list:
            fileFullPath = Path(item)
            fileName = fileFullPath.name
            print(f"Processing {fileName}")
            if fileName.endswith(".docx"):
                try:
                    document = DocxDocument(fileFullPath)
                    content = "\n".join([paragraph.text for paragraph in document.paragraphs])
                except Exception as err:
                    logger.error(f"Error processing document: {err}")
            elif fileName.endswith(".pdf"):
                try:
                    content = extractPDFContent(fileFullPath)
                except Exception as err:
                    logger.error(f"Error processing document: {err}")
            elif fileName.endswith(".pptx"):
                try:
                    content = extractPPTXContent(fileFullPath)
                except Exception as err:
                    logger.error(f"Error processing document: {err}")
            file_info = {
                'id': str(uuid.uuid4()),
                'name': fileName,
                'created_datetime': datetime.fromtimestamp(fileFullPath.stat().st_ctime),
                'created_by': str(fileFullPath.stat().st_uid),
                'size': fileFullPath.stat().st_size,
                'last_modified_datetime': datetime.fromtimestamp(fileFullPath.stat().st_mtime),
                'last_modified_by': str(fileFullPath.stat().st_uid),
                'source': item,
                # Israel
                "parentObject": f"{file_system_path}",
                ########
                'typedef': 'FileSystem',
                'content': content,
            }
            files_to_process.append(file_info)

    return files_to_process

def sharepointFileSampleList(input_file_list, file_sample_size):
    """
    Generates a list of files to process from an input list of files.

    Args:
        input_file_list (list): A list of dictionaries containing file information.
            Each dictionary should have a 'name' key representing the file name.
        file_sample_size (int): The desired sample size of files to process.

    Returns:
        list: A list of file names to process. If the number of files found is greater
        than the sample size, a random sample of files is returned. Otherwise, all files
        in the input list are included.

    Example:
        input_files = [{'name': 'file1.txt'}, {'name': 'file2.txt'}, {'name': 'file3.txt'}]
        sample_size = 2
        result = sharepointFileSampleList(input_files, sample_size)
        # Example output: ['file1.txt', 'file3.txt']
    """
    num_files_found = len(input_file_list)
    files_to_process = []

    if num_files_found > file_sample_size:
        # Randomly sample files if more files are found than the desired sample size
        for ind in random.sample(range(num_files_found), file_sample_size):
            files_to_process.append(input_file_list[ind]['name'])
    else:
        # Include all files if the number of files found is less than or equal to the sample size
        for ind in input_file_list:
            files_to_process.append(ind['name'])

    return files_to_process

def listFilesystemFiles(mypath, extensions):
    """
    Lists files in a directory matching specified extensions.

    Args:
        mypath (str): Path to the directory.
        extensions (list): List of file extensions (e.g., ['.txt', '.csv', '.jpg']).

    Returns:
        list: List of filenames matching the specified extensions.
    """
    return [os.path.join(dirpath, f) for (dirpath, _, filenames) in os.walk(mypath) for f in filenames if any(f.endswith(ext) for ext in extensions)]

def extractPDFContent(fileFullPath):
    # Function to convert PDF files to Text
    try:
        pdffileobj = open(fileFullPath,'rb')
        pdfreader = PyPDF2.PdfReader(pdffileobj)
        num_pages = len(pdfreader.pages)
        content = ''
        for page in range(num_pages):
            pageobj = pdfreader.pages[page]
            content += ''.join(pageobj.extract_text())
        return content
    except Exception as err:
        logger.error(f"Error processing document: {err}")
        return None

def extractPPTXContent(fileFullPath):
    # Function to convert PPTX files to Text
    try:
        presentation = Presentation(fileFullPath)
        content = ''
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    content += shape.text + '\n'
        return content
    except Exception as err:
        logger.error(f"Error processing document: {err}")
        return None

def getAADToken(tenant_id: str, client_id: str, client_secret: str, resource_url: str):
    """
    Authenticates Service Principal to the provided Resource URL, and returns the OAuth Access Token
    """
    url = f"https://login.microsoftonline.com/{tenant_id}/oauth2/token"
    payload= f'grant_type=client_credentials&client_id={client_id}&client_secret={client_secret}&resource={resource_url}'
    headers = {
    'Content-Type': 'application/x-www-form-urlencoded'
    }
    response = requests.request("POST", url, headers=headers, data=payload)
    access_token = json.loads(response.text)['access_token']
    return access_token

def moveCollection(collectionId: str, endpoint: str, token: str, guids):
    url = f"{endpoint}/datamap/api/entity/moveTo?api-version=2023-09-01&collectionId={collectionId}"
    payload={
        "entityGuids": guids
    }
    headers = {
        'Authorization': f'Bearer {token}',
        'Content-Type': 'application/json'
    }
    payload_json = json.dumps(payload)
    response = requests.request("POST", url, headers=headers, data=payload_json)
    return json.loads(response.text)

def listAssets(endpoint: str, token: str, query: str):
    url = f"{endpoint}/datamap/api/search/query?api-version=2023-09-01"
    payload={
        "keywords": query
    }
    headers = {
        'Authorization': f'Bearer {token}',
        'Content-Type': 'application/json'
    }
    payload_json = json.dumps(payload)
    response = requests.request("POST", url, headers=headers, data=payload_json)
    return json.loads(response.text)

def deleteAssets(purviewClient, endpoint: str, token: str, query: str):
    url = f"{endpoint}/datamap/api/search/query?api-version=2023-09-01"
    payload={
        "keywords": query
    }
    headers = {
        'Authorization': f'Bearer {token}',
        'Content-Type': 'application/json'
    }
    payload_json = json.dumps(payload)
    response = requests.request("POST", url, headers=headers, data=payload_json)
    assets = [asset['id'] for asset in json.loads(response.text)['value']]
    delResponse = purviewClient.delete_entity(guid=assets)
    return json.dumps(delResponse, indent=2)

def estimateTokens(fileContents,textLength,classificationsStr,aiModel):
    """
    Estimates the total number of tokens in a given user input, including system messages and user prompts.

    Args:
        fileContents (List[Doc]): List of document contents to be analyzed.
        textLenght (int): Number of characters to be analyzed from of each document.
        classificationsStr (str): String containing the list of classifications to be evaluated against each document.

    Returns:
        totalTokens: Int containing the estimated total number of tokens.
    """
    totalTokens = 0
    for ind in range(len(fileContents)):
        fileContent = fileContents[ind]
        # Text extracted from scanned documents up to the number of characters specified by textLength
        text = fileContent['content'][:textLength]
        if re.search(r'[a-zA-Z0-9]',text):
            # Generate user prompt
            userPrompt = (
                "role: system, content: You are an AI assistant that helps people classify information contained in documents."
                +"role: user, content: "
                +"Classify the following text: "+text
                +"Using only any of the following classifications: \n"+classificationsStr
                +"The answer should only return one of the listed classifications."
            )
            # enc = tiktoken.get_encoding("cl100k_base")
            # enc = tiktoken.encoding_for_model('gpt-4')
            enc = tiktoken.encoding_for_model(aiModel)
            numTokens = len(enc.encode(userPrompt))
        else:
            numTokens = 0
        totalTokens += numTokens
    return totalTokens

def unstructuredDataClassification(fileContents,textLength,llmClient,llmodel,classificationsStr):
    """
    This function will run the contents of each file contained in the subset against the 
    Azure Open AI large language model (GPT-4). The result will be a list of dictionaries
    containing file metadata and also the classification assigned by the LLM model.
    """
    for ind in range(len(fileContents)):
        fileContent = fileContents[ind]
        # Text extracted from scanned documents up to the number of characters specified by textLength
        text = fileContent['content'][:textLength]
        if re.search(r'[a-zA-Z0-9]',text):
            # Generate user prompt
            userPrompt = (
                            "Classify the following text: "+text
                            +"Using only any of the following classifications:\n"+classificationsStr
                            +"The answer should only return one of the listed classifications."
                        )
            # Submit prompts to LLM model
            response = llmClient.complete(
                messages=[
                    SystemMessage(content="You are an insurance specialist that helps people classify information contained in documents."),
                    UserMessage(content=userPrompt),
                ],
                max_tokens=4096,
                temperature=1.0,
                top_p=1.0,
                model=llmodel
            )
            # Original code for Azure Open AI
            # response = llmClient.chat.completions.create(
            #     model=azureOpenAIDeploymentName,
            #     messages=[
            #         {"role": "system", "content": "You are an AI assistant that helps people classify information contained in documents."},
            #         {"role": "user", "content": userPrompt},
            #     ]
            # )
            print(f"Processing {fileContent['name']} -> {response.choices[0].message.content}")
            fileContents[ind].update({'classification':response.choices[0].message.content})
            # Comment out the line below to determine the exact number of tokens consumed by each document
            print(response.usage)
        else:
            print(f"Processing {fileContent['name']} -> Empty Content")
            fileContents[ind].update({'classification':'Empty Content'})
    return fileContents

def loadPurviewAssets(purviewClient,allFileContent):
    """
    Load parent level assets for SharePoint and FileSystem data sources.
    """
    batchEntities = []
    newGuid = -1000
    # classificationsFound = list(set([file['classification'] for file in allFileContent]))
    # classificationsFound = list(set(classificationsFound))
    file = allFileContent[0]
    print(f"Processing {file['parentObject']} - {file['classification']}")

    newEntity = AtlasEntity(
        name=file["parentObject"],
        typeName=file["typedef"],
        qualified_name=f"customScanner://{file['parentObject']}",
        guid=newGuid
    )

    batchEntities.append(newEntity)

    # Convert the individual entities into json before uploading.
    results = purviewClient.upload_entities(
        batch = batchEntities,
        batch_size=20
    )

    # # Get the Guids for us to work with
    guids = [v for v in results["guidAssignments"].values()]
    return guids

def rollupClassifications(allFileContent):
    """
    Obtain parent level classifications to represent the classifications of the individual
    files stored in the data source.
    """
    classificationsFound = list(set([file['classification'] for file in allFileContent]))
    return classificationsFound

def applyPurviewClassifications(purviewClient,guids,classificationsFound):
    """
    Apply Classifications to Purview parent resources
    """
    atlasClassifications = [AtlasClassification(classification_name) for classification_name in classificationsFound]
    results = purviewClient.classify_entity(
        guid=guids[0], 
        classifications=atlasClassifications,
        force_update=True
    )
    return results